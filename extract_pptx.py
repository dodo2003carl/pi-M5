import sys
import os

pptx_path = "Presentación del PI M5.pptx"

try:
    from pptx import Presentation
    print("Using python-pptx")
    prs = Presentation(pptx_path)
    text_content = []
    
    for slide_number, slide in enumerate(prs.slides):
        text_content.append(f"--- Slide {slide_number + 1} ---")
        for shape in slide.shapes:
            if hasattr(shape, "text"):
                text_content.append(shape.text)
    
    full_text = "\n".join(text_content)
    
    with open("pptx_content.txt", "w", encoding="utf-8") as f:
        f.write(full_text)
    print("PPTX content written to pptx_content.txt")

except ImportError:
    print("No python-pptx library found. Please install python-pptx.")
except Exception as e:
    print(f"Error processing PPTX: {e}")
